import sys
import struct
from pptx import Presentation
from pptx.util import Inches, Px
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE

try:

# the first argument is the name of this script. Useless.
    list_of_img_paths = sys.argv[2:]

    prs = Presentation()

    for img_path in list_of_img_paths:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)


## add rectangle to make background back. ###
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(12))

# set fill type to solid color first
#    print 1
        shape.fill.solid()

# set foreground (fill) color to a specific RGB color
#    print 2
        shape.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
        #shape.fill.fore_color.rgb = RGBColor(0xFB, 0x8F, 0x00)

# change to a theme color
#    print 3
        shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1

# set lighter or darker, -0.2 is 20% darker, 0.4 is 40% lighter
#    print 4
        shape.fill.fore_color.brightness = -0.9
## done with add rectangle. ###



## scaling image
        im = open(img_path, 'rb').read()
        i = im.find('\xff\xc0') + 5
        original_height, original_width = struct.unpack('>HH', im[i:i+4])

        height_in_inch = 6.0
        width_in_inch  = 10.0
        factor_h = height_in_inch / original_height
        factor_w = width_in_inch  / original_width

        if factor_w < factor_h:
            width = Inches(original_width * factor_w)
            height = Inches(original_height * factor_w)
            left = Inches(0)
            top = Inches((8.0 - (original_height * factor_w)) / 2)
        else:
            width = Inches(original_width * factor_h)
            height = Inches(original_height * factor_h)
            left = Inches((10.0 - (original_width * factor_h)) / 2)
            top  = Inches(1)
### add a picture

        pic = slide.shapes.add_picture(img_path, left, top, width, height)

    prs.save(sys.argv[1])
except:
    print 'Error occurred'
    print sys.exc_info()[0]
