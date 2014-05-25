# Updated version: Add pictures and text

# python make_ppt.py <file_name>.pptx <path1> "<comment1>" <path2> "<comment2>" ...

import sys
import struct
from pptx import Presentation
from pptx.util import Inches, Px
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

try:

# the first argument is the name of this script. Useless.
    my_args = sys.argv[2:]

    prs = Presentation()

    index = 0
    while index < len(my_args):
        img_path = my_args[index]
        index += 1
        my_text     = my_args[index]
        index += 1
#    print "7373 text", my_text

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        height_in_inch = 6.0
        width_in_inch  = 10.0
        screen_height = 9.0
        screen_width = 10.0

## add rectangle to make background back. ###
## To make the whole background black, change the last two argumets to 10, 12
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(width_in_inch), Inches(height_in_inch))

# set fill type to solid color first
#    print 1
        shape.fill.solid()

# set foreground (fill) color to a specific RGB color
#    print 2
        shape.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
        #shape.fill.fore_color.rgb = RGBColor(0xFB, 0x8F, 0x00)

# change to a theme color
#    print 3
        shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1

# set lighter or darker, -0.2 is 20% darker, 0.4 is 40% lighter
#    print 4
        shape.fill.fore_color.brightness = -0.9
## done with add rectangle. ###



## scaling image
        im = open(img_path, 'rb').read()
        i = im.find('\xff\xc0') + 5
        original_height, original_width = struct.unpack('>HH', im[i:i+4])

        factor_h = height_in_inch / original_height
        factor_w = width_in_inch  / original_width

        if factor_w < factor_h:
            width = Inches(original_width * factor_w)
            height = Inches(original_height * factor_w)
            left = Inches(0)
            top = Inches((height_in_inch - (original_height * factor_w)) / 2)
        else:
            width = Inches(original_width * factor_h)
            height = Inches(original_height * factor_h)
            left = Inches((screen_width - (original_width * factor_h)) / 2)
            top  = Inches(0)
### add a picture

#    print "7373 image path ---" + img_path + "-----"
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
#    pic = slide.shapes.add_picture("baseball.jpg", left, top, width, height)


#######
##    slide = prs.slides.add_slide(blank_slide_layout)
##    title = slide.shapes.title
##    subtitle = slide.placeholders[1]
##
##    title.text = "Hello, World!"
##    subtitle.text = "python-pptx was here!"
#    txBox = slide.shapes.add_textbox(left, top, width, height)
#    txBox = slide.shapes.add_textbox(0, screen_height_for_picture, width_in_inch, screen_height_for_picture - height_in_inch)

        # delta is the distance between text and picture
        delta = 0.5
        txBox = slide.shapes.add_textbox(Inches(0), Inches(height_in_inch + delta), Inches(width_in_inch), Inches(screen_height - height_in_inch))
        tf = txBox.textframe

################### changing font
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = my_text

        font = run.font
    #    font.name = 'Calibri'
        font.size = Pt(30)
#        font.bold = True
################### end changing font

#        tf.text = my_text
#######


    prs.save(sys.argv[1])
except:
    print 'Error occurred'
    print sys.exc_info()[0]
